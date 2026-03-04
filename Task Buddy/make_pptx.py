from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette ───────────────────────────────────────────────
NAVY       = RGBColor(0x0E, 0x1F, 0x3D)
ORANGE     = RGBColor(0xFF, 0x69, 0x00)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF4, 0xF6, 0xF9)
MID_GRAY   = RGBColor(0x6B, 0x7A, 0x99)
SOFT_BLUE  = RGBColor(0xD6, 0xE4, 0xFF)
GREEN      = RGBColor(0x00, 0xA8, 0x6B)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── Helpers ───────────────────────────────────────────────
def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)

def fill_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, fill_color, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    return shape

def add_text(slide, text, l, t, w, h, size, bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return tb

def add_label(slide, text, l, t, w, h):
    """Orange ALL-CAPS label"""
    add_text(slide, text, l, t, w, h, 9, bold=True, color=ORANGE, align=PP_ALIGN.LEFT)

def add_bullet_box(slide, items, l, t, w, h, bg_color, text_color, title=None):
    add_rect(slide, l, t, w, h, bg_color)
    if title:
        add_text(slide, title, l+0.15, t+0.12, w-0.3, 0.35, 10, bold=True, color=ORANGE)
        start_y = t + 0.45
    else:
        start_y = t + 0.15
    line_h = (h - (0.45 if title else 0.3)) / max(len(items), 1)
    for i, item in enumerate(items):
        add_text(slide, f"• {item}", l+0.15, start_y + i*line_h, w-0.3, line_h+0.05, 11, color=text_color)

# ══════════════════════════════════════════════════════════
#  SLIDE 1 – Cover
# ══════════════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl, NAVY)

# Left accent bar
add_rect(sl, 0, 0, 0.12, 7.5, ORANGE)

# Big decorative circle (top-right)
circ = sl.shapes.add_shape(9, Inches(9.8), Inches(-1.2), Inches(4.5), Inches(4.5))
circ.fill.solid()
circ.fill.fore_color.rgb = RGBColor(0x1A, 0x3A, 0x6B)
circ.line.fill.background()

# Small orange dot
dot = sl.shapes.add_shape(9, Inches(11.2), Inches(5.5), Inches(0.5), Inches(0.5))
dot.fill.solid()
dot.fill.fore_color.rgb = ORANGE
dot.line.fill.background()

# Tag line
add_label(sl, "PROJECT OVERVIEW  ·  HUMAN-CENTRED DESIGN", 0.4, 0.4, 10, 0.4)

# Title
add_text(sl, "Thinking Before\nBuilding", 0.4, 1.1, 9, 2.6, 56, bold=True, color=WHITE)

# Sub-line under title
add_text(sl,
    "How intentional design — not just code — creates products\n"
    "that truly serve real people.",
    0.4, 3.55, 8.2, 1.2, 16, color=RGBColor(0x8F, 0xA8, 0xCC))

# Three pill badges
badge_texts = ["📋 Project Description", "👤 User Personas", "🔀 User Flows"]
for i, bt in enumerate(badge_texts):
    bx = 0.4 + i * 3.1
    rect = sl.shapes.add_shape(5, Inches(bx), Inches(5.15), Inches(2.8), Inches(0.48))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0x1E, 0x38, 0x6B)
    rect.line.color.rgb = ORANGE
    rect.line.width = Pt(1)
    add_text(sl, bt, bx+0.1, 5.18, 2.6, 0.4, 10.5, bold=True, color=WHITE)

# Bottom strip
add_rect(sl, 0, 7.05, 13.33, 0.45, RGBColor(0x06, 0x12, 0x26))
add_text(sl, "Start early  ·  Be intentional  ·  Take ownership",
         0, 7.1, 13.33, 0.35, 10, color=MID_GRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════
#  SLIDE 2 – What Is This Project About?
# ══════════════════════════════════════════════════════════
sl2 = blank_slide(prs)
fill_bg(sl2, LIGHT_GRAY)
add_rect(sl2, 0, 0, 13.33, 1.35, NAVY)
add_label(sl2, "SLIDE 1 OF 4  ·  THE TOPIC", 0.45, 0.17, 8, 0.35)
add_text(sl2, "What Is This Project About?", 0.45, 0.48, 12, 0.8, 34, bold=True, color=WHITE)

# Intro paragraph
add_text(sl2,
    "Every day, people struggle with tools that are confusing, hard to use, or simply not built for them.\n"
    "This project is about changing that — by designing digital products that put PEOPLE first, before any line of code is written.",
    0.45, 1.55, 12.4, 1.1, 14, color=RGBColor(0x2D, 0x3A, 0x55))

# Three concept cards
card_data = [
    ("📋", "Clear Project\nDescription",
     "Before building anything, you must be able to explain exactly WHAT you are making, WHO it is for, and WHY it matters — in plain language anyone can understand."),
    ("👤", "User Personas",
     "A 'persona' is a portrait of a real type of person who will use your product. We give them a name, a story, and real frustrations — so every design decision focuses on helping them."),
    ("🔀", "User Flows",
     "A user flow is a simple map that shows the steps a person takes to complete a task — like signing up or creating something. It helps us spot confusion before anything is built."),
]

card_w = 3.8
gap = 0.37
for i, (icon, title, body) in enumerate(card_data):
    cx = 0.42 + i * (card_w + gap)
    add_rect(sl2, cx, 3.0, card_w, 3.75, WHITE)
    # top accent bar
    add_rect(sl2, cx, 3.0, card_w, 0.1, ORANGE)
    add_text(sl2, icon, cx+0.2, 3.12, 0.6, 0.55, 26)
    add_text(sl2, title, cx+0.2, 3.68, card_w-0.4, 0.65, 14, bold=True, color=NAVY)
    add_text(sl2, body, cx+0.2, 4.35, card_w-0.4, 2.2, 11.5, color=RGBColor(0x4A, 0x55, 0x68))

# Bottom quote
add_rect(sl2, 0, 6.9, 13.33, 0.6, NAVY)
add_text(sl2, '"This project is not only about coding — it is about thinking, designing, and building responsibly."',
         0.4, 6.95, 13, 0.5, 10.5, italic=True, color=RGBColor(0x8F, 0xA8, 0xCC), align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════
#  SLIDE 3 – The Challenge
# ══════════════════════════════════════════════════════════
sl3 = blank_slide(prs)
fill_bg(sl3, WHITE)
add_rect(sl3, 0, 0, 13.33, 1.35, NAVY)
add_label(sl3, "SLIDE 2 OF 4  ·  THE CHALLENGE", 0.45, 0.17, 8, 0.35)
add_text(sl3, "What Problem Are We Solving?", 0.45, 0.48, 12, 0.8, 34, bold=True, color=WHITE)

# Left column – problem
add_rect(sl3, 0.35, 1.5, 5.9, 5.2, RGBColor(0xFF, 0xF4, 0xED))
add_rect(sl3, 0.35, 1.5, 5.9, 0.08, ORANGE)
add_text(sl3, "😟  The Reality", 0.55, 1.6, 5.5, 0.5, 15, bold=True, color=NAVY)
problems = [
    "Developers often start coding without a clear plan",
    "Products get built for the developer, not the user",
    "Teams skip thinking about WHO will actually use the product",
    "Features are added without understanding real needs",
    "Users end up confused, frustrated, or simply stop using the product",
]
for i, p in enumerate(problems):
    add_text(sl3, f"✗  {p}", 0.55, 2.18 + i*0.82, 5.5, 0.75, 12, color=RGBColor(0x7A, 0x3F, 0x1A))

# Right column – consequences
add_rect(sl3, 6.85, 1.5, 6.1, 5.2, RGBColor(0xEA, 0xF0, 0xFF))
add_rect(sl3, 6.85, 1.5, 6.1, 0.08, RGBColor(0x3C, 0x67, 0xEA))
add_text(sl3, "⚠️  The Impact of Skipping Design", 7.05, 1.6, 5.7, 0.5, 15, bold=True, color=NAVY)
impacts = [
    "Time wasted rebuilding features that don't work for users",
    "Products that look finished but solve the wrong problem",
    "Poor user experience = low adoption, bad reviews",
    "Missed opportunity to create something truly meaningful",
]
for i, imp in enumerate(impacts):
    add_text(sl3, f"→  {imp}", 7.05, 2.18 + i*1.0, 5.7, 0.9, 12.5, color=RGBColor(0x1E, 0x3A, 0x7A))

# ══════════════════════════════════════════════════════════
#  SLIDE 4 – The Solution & User Personas
# ══════════════════════════════════════════════════════════
sl4 = blank_slide(prs)
fill_bg(sl4, LIGHT_GRAY)
add_rect(sl4, 0, 0, 13.33, 1.35, NAVY)
add_label(sl4, "SLIDE 3 OF 4  ·  THE SOLUTION", 0.45, 0.17, 8, 0.35)
add_text(sl4, "The Solution: Design With People in Mind", 0.45, 0.48, 12, 0.8, 30, bold=True, color=WHITE)

# Top row – 3 step approach
steps = [
    ("01", "Write a Clear\nProject Description",
     "Put into plain words: What does this do? Who is it for? What problem does it solve? This forces clarity before building begins."),
    ("02", "Create User\nPersonas",
     "Give your users a face and a story. Understand their age, goals, frustrations, and daily habits — so every feature is built FOR them."),
    ("03", "Map the\nUser Journey",
     "Draw simple step-by-step paths that show how a person reaches their goal. Find the confusing parts and fix them on paper — not after."),
]
sw = 3.9
for i, (num, title, body) in enumerate(steps):
    cx = 0.35 + i*(sw+0.27)
    add_rect(sl4, cx, 1.55, sw, 2.55, WHITE)
    add_rect(sl4, cx, 1.55, sw, 0.08, ORANGE)
    add_text(sl4, num, cx+0.18, 1.65, 0.9, 0.65, 28, bold=True, color=ORANGE)
    add_text(sl4, title, cx+0.18, 2.28, sw-0.35, 0.6, 14, bold=True, color=NAVY)
    add_text(sl4, body, cx+0.18, 2.92, sw-0.35, 1.1, 11.5, color=RGBColor(0x4A, 0x55, 0x68))

# Persona showcase
add_text(sl4, "Meet the Users", 0.35, 4.35, 6, 0.45, 17, bold=True, color=NAVY)
personas = [
    ("👩‍🎓", "Amara, 21", "University Student",
     ["Overwhelmed by assignments & deadlines", "Needs a simple way to track her tasks", "Uses her phone for everything"]),
    ("👨‍💼", "Daniel, 34", "Freelance Designer",
     ["Juggles multiple client projects", "Loses track of what is done vs. pending", "Wants clarity, not complexity"]),
    ("👩‍💻", "Kalisa, 23", "Recent Graduate",
     ["Starting a new job, learning new workflows", "Needs to build good habits from day one", "Values clean, intuitive tools"]),
]
pw = 3.9
for i, (icon, name, role, traits) in enumerate(personas):
    px = 0.35 + i*(pw+0.27)
    add_rect(sl4, px, 4.82, pw, 2.42, WHITE)
    add_rect(sl4, px, 4.82, pw, 0.07, GREEN)
    add_text(sl4, icon, px+0.18, 4.9, 0.7, 0.55, 22)
    add_text(sl4, name, px+0.95, 4.9, pw-1.1, 0.35, 13, bold=True, color=NAVY)
    add_text(sl4, role, px+0.95, 5.24, pw-1.1, 0.3, 10.5, color=MID_GRAY)
    for j, t in enumerate(traits):
        add_text(sl4, f"· {t}", px+0.18, 5.6+j*0.52, pw-0.35, 0.48, 10.5, color=RGBColor(0x4A, 0x55, 0x68))

# ══════════════════════════════════════════════════════════
#  SLIDE 5 – Impact & Call to Action
# ══════════════════════════════════════════════════════════
sl5 = blank_slide(prs)
fill_bg(sl5, NAVY)
add_rect(sl5, 0, 0, 0.12, 7.5, ORANGE)

# Decorative bottom circle
circ2 = sl5.shapes.add_shape(9, Inches(9.2), Inches(4.5), Inches(5.5), Inches(5.5))
circ2.fill.solid()
circ2.fill.fore_color.rgb = RGBColor(0x1A, 0x3A, 0x6B)
circ2.line.fill.background()

add_label(sl5, "SLIDE 4 OF 4  ·  THE IMPACT", 0.4, 0.28, 9, 0.35)
add_text(sl5, "What Changes When\nWe Design Responsibly?", 0.4, 0.65, 9.5, 1.8, 38, bold=True, color=WHITE)

impacts_pos = [
    ("🎯", "Products people actually want to use",
     "When you design around real needs, adoption is natural — not forced."),
    ("⏱️", "Less time wasted, more value delivered",
     "Catching problems in the design stage costs nothing. Fixing them after launch costs everything."),
    ("❤️", "Technology that earns trust",
     "People trust tools that feel like they were made FOR them — thoughtfully, not carelessly."),
    ("🚀", "A stronger foundation for building",
     "A clear description + personas + user flows = a roadmap everyone understands, not just developers."),
]

for i, (icon, title, body) in enumerate(impacts_pos):
    row = i // 2
    col = i % 2
    ix = 0.4 + col * 6.3
    iy = 2.75 + row * 2.1
    add_rect(sl5, ix, iy, 5.9, 1.75, RGBColor(0x12, 0x28, 0x55))
    add_rect(sl5, ix, iy, 5.9, 0.07, ORANGE)
    add_text(sl5, icon, ix+0.15, iy+0.12, 0.65, 0.55, 24)
    add_text(sl5, title, ix+0.85, iy+0.12, 4.85, 0.45, 13, bold=True, color=WHITE)
    add_text(sl5, body, ix+0.85, iy+0.58, 4.85, 0.95, 11.5, color=RGBColor(0x8F, 0xA8, 0xCC))

# Bottom strip
add_rect(sl5, 0, 7.05, 13.33, 0.45, RGBColor(0x06, 0x12, 0x26))
add_text(sl5, "Start early  ·  Be intentional  ·  Take ownership",
         0, 7.1, 13.33, 0.35, 10, color=MID_GRAY, align=PP_ALIGN.CENTER)

# ── Save ──────────────────────────────────────────────────
out = r"c:\Users\07892\OneDrive\Desktop\kalisa\my-project\HumanCentredDesign.pptx"
prs.save(out)
print(f"Saved → {out}")
